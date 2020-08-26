# -*- coding:utf-8 -*-

import csv
from collections import defaultdict
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from bible_api import BibleAPI

BLANK_SLIDE = 6

RR_FONT_SIZE = 68
BIBLE_FONT_SIZE = 80
BIBLE_SPACE = 0
TITLE_FONT_SIZE = 75
TITLE_SPACE = 330
SUBTITLE_FONT_SIZE = 45
SUBTITLE_SPACE = 410

WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(249, 166, 2)


class Bible2ppt:
    def __init__(self):
        self.api = BibleAPI()
        self.ppt = Presentation('./templates/base_template.pptx')

    def __generate_slide(self, template):
        """
        해당 템플릿 사진을 배경으로 하는 새로운 슬라이드를 반환한다.

        :params:
            * template (str) : 원하는 템플릿 이름

        :returns:
            * slide (slide) : 해당 템플릿 사진을 배경으로 하는 슬라이드
        """
        left = top = 0
        slide_layout = self.ppt.slide_layouts[BLANK_SLIDE]
        slide = self.ppt.slides.add_slide(slide_layout)

        background = slide.shapes.add_picture(f'templates/{template}.png', left, top, height=self.ppt.slide_height)
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)

        return slide
    
    def __add_text(self, slide, text, space, size, color, is_center):
        """
        해당 슬라이드에 텍스트를 적는다.

        :params:
            * slide (slide) : 슬라이드
            * text (str) : 슬라이드에 적을 텍스트
            * space (int) : 텍스트 상자의 y 좌표
            * size (int) : 텍스트 크기
            * color (RGBColor) : 텍스트 색상
            * is_center (bool) : 텍스트 중앙 정렬 여부
        """
        margin = Inches(0.1)
        tb = slide.shapes.add_textbox(margin, margin, self.ppt.slide_width-2*margin, margin)
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = text
        p.font.bold = True
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_before = Pt(space)
        
        if is_center:
            p.alignment = PP_ALIGN.CENTER

    def __set_responsive_reading(self, index):
        """
        '성시교독'과 해당 성시 교독이 적힌 슬라이드를 추가한다.

        :params:
            * index (int) : 해당 성시 교독 번호
        """
        title, contents = self.api.get_responsive_reading(index)
        cover = self.__generate_slide('rr')
        self.__add_text(cover, title, TITLE_SPACE, TITLE_FONT_SIZE, WHITE, True)

        for idx, content in enumerate(contents):
            color = WHITE if not idx % 2 else YELLOW
            space = 0 if not idx % 2 else 235

            if not idx % 2:
                slide = self.__generate_slide('background')
            self.__add_text(slide, content, space, RR_FONT_SIZE, color, False)

    def __set_bible(self, lesson, assigner):
        """
        '말씀봉독'과 해당 성경 말씀이 적힌 슬라이드를 추가한다.

        사회자는 흰색, 회동은 노란색 부분을 봉독한다.

        :params:
            * lesson (str) : 성경 챕터, 시작 장, 시작 절, 끝 장, 끝 절
            * assigner (str) : 성경 봉독 담당자
        """
        chapter, aa, bb, cc, dd = lesson.split(',')
        title, contents = self.api.get_bible(chapter, aa, bb, cc, dd)

        cover = self.__generate_slide('bible')
        self.__add_text(cover, title, TITLE_SPACE, TITLE_FONT_SIZE, WHITE, True)
        self.__add_text(cover, assigner, SUBTITLE_SPACE, SUBTITLE_FONT_SIZE, WHITE, True)

        for idx, content in enumerate(contents):
            slide = self.__generate_slide('background')
            text = f'{content.get("section")} - {content.get("sentence")}'
            self.__add_text(slide, text, BIBLE_SPACE, BIBLE_FONT_SIZE, WHITE, False)
            
    def __set_offering(self, assigner):
        """
        '봉헌', '봉헌찬양', '봉헌기도' 슬라이드를 추가한다.

        :params:
            * assinger (str) : 봉헌 찬양 담당자
        """
        self.__generate_slide('offering')
        
        slide = self.__generate_slide('offering_praise')
        self.__add_text(slide, assigner, TITLE_SPACE, TITLE_FONT_SIZE, WHITE, True)

        self.__generate_slide('offering_prayer')

    def __set_sermon(self, info):
        """
        '설교' 슬라이드를 추가한다.

        :params:
            * info (string) : 설교제목, 담당자
        """
        title, assigner = info.split(',')
        cover = self.__generate_slide('sermon')
        self.__add_text(cover, title, TITLE_SPACE, TITLE_FONT_SIZE, WHITE, True)
        self.__add_text(cover, assigner, SUBTITLE_SPACE, SUBTITLE_FONT_SIZE, WHITE, True)

    def __get_worship_info(self):
        """
        worship_info.csv를 읽어서 예배 정보를 반환한다.

        :return:
            * info (dict) : 예배 정보
        """
        info = dict()
        with open('worship_info.csv', encoding='UTF-8') as f:
            rows = csv.reader(f, delimiter='|')
            for r in rows:
                info[r[0]] = r[1]

        return info

    def run(self):
        """
        예배 순서에 맞게 ppt를 자동으로 생성한 뒤, 오늘 날짜를 파일명으로 하여 저장한다.
        """
        info = self.__get_worship_info()

        self.__set_responsive_reading(info.get('rr'))
        self.__generate_slide('hymn')
        
        slide = self.__generate_slide('prayer')
        self.__add_text(slide, info.get('prayer'), TITLE_SPACE, TITLE_FONT_SIZE, WHITE, True)
  
        self.__set_offering(info.get('offering'))
        self.__set_bible(info.get('bible'), info.get('reader'))

        self.__generate_slide('hymn_noel')
        self.__set_sermon(info.get('sermon'))
        self.__generate_slide('altar_call')
        self.__generate_slide('benediction')

        today = datetime.now().strftime('%Y-%m-%d')
        self.ppt.save(f'{today}.pptx')


if __name__ == '__main__':
    b2p = Bible2ppt()
    b2p.run()